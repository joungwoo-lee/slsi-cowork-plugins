"""Convert mail body + attachment formats into a single unified markdown."""
from __future__ import annotations

import io
import logging
from pathlib import Path

from markdownify import markdownify as html_to_md

log = logging.getLogger(__name__)

ATTACHMENT_HEADER = "\n\n---\n\n[첨부파일: {name}]\n\n"

# Suffixes treated as plain-text-with-codefence. Mapping value is the codefence
# language hint; "" means embed as raw text without a fence.
_TEXT_SUFFIXES: dict[str, str] = {
    ".txt": "",
    ".log": "",
    ".md": "",
    ".csv": "csv",
    ".tsv": "tsv",
    ".json": "json",
    ".xml": "xml",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".ini": "ini",
    ".conf": "",
    ".cfg": "",
    ".sql": "sql",
    ".py": "python",
    ".js": "javascript",
    ".ts": "typescript",
    ".sh": "bash",
    ".bat": "bat",
    ".ps1": "powershell",
}

_IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg"}


def body_to_markdown(html: str, plain: str, max_chars: int) -> str:
    if html and html.strip():
        try:
            md = html_to_md(html, heading_style="ATX")
        except Exception as exc:  # pragma: no cover
            log.warning("markdownify failed, falling back to plain: %s", exc)
            md = plain or ""
    else:
        md = plain or ""
    return md.strip()[:max_chars]


# -----------------------------------------------------------------------------
# Per-format converters
# -----------------------------------------------------------------------------


def pdf_to_markdown(data: bytes, max_chars: int) -> str:
    import fitz  # pymupdf

    parts: list[str] = []
    used = 0
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page_num, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if not text:
                continue
            chunk = f"### 페이지 {page_num}\n\n{text}\n"
            if used + len(chunk) > max_chars:
                parts.append(chunk[: max_chars - used])
                break
            parts.append(chunk)
            used += len(chunk)
    return "\n".join(parts)


def docx_to_markdown(data: bytes, max_chars: int) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    out: list[str] = []
    used = 0

    def push(line: str) -> bool:
        nonlocal used
        if used + len(line) + 1 > max_chars:
            out.append(line[: max_chars - used])
            return False
        out.append(line)
        used += len(line) + 1
        return True

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "2"
            line = f"{'#' * min(int(level), 6)} {text}"
        else:
            line = text
        if not push(line):
            return "\n".join(out)

    for tbl_idx, table in enumerate(doc.tables, start=1):
        if not push(f"\n**표 {tbl_idx}**\n"):
            return "\n".join(out)
        for row in table.rows:
            cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
            if not push("| " + " | ".join(cells) + " |"):
                return "\n".join(out)
    return "\n".join(out)


def xlsx_to_markdown(data: bytes, max_chars: int) -> str:
    """Excel — render each sheet as a pipe table. Empty sheets skipped."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out: list[str] = []
    used = 0

    def push(line: str) -> bool:
        nonlocal used
        chunk = line + "\n"
        if used + len(chunk) > max_chars:
            return False
        out.append(chunk)
        used += len(chunk)
        return True

    for ws in wb.worksheets:
        if not push(f"### 시트: {ws.title}"):
            return "".join(out)
        any_row = False
        for row in ws.iter_rows(values_only=True):
            cells = [
                "" if c is None else str(c).replace("\n", " ").replace("|", "/").strip()
                for c in row
            ]
            if not any(cells):
                continue
            any_row = True
            if not push("| " + " | ".join(cells) + " |"):
                return "".join(out)
        if not any_row:
            push("(빈 시트)")
        push("")
    return "".join(out).strip()


def pptx_to_markdown(data: bytes, max_chars: int) -> str:
    """PowerPoint — slide-by-slide text including text frames, tables, notes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    used = 0

    def push(line: str) -> bool:
        nonlocal used
        chunk = line + "\n"
        if used + len(chunk) > max_chars:
            return False
        out.append(chunk)
        used += len(chunk)
        return True

    for i, slide in enumerate(prs.slides, start=1):
        if not push(f"### 슬라이드 {i}"):
            return "".join(out)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text or "" for run in para.runs).strip()
                    if text and not push(text):
                        return "".join(out)
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    cells = [
                        (cell.text or "").strip().replace("\n", " ").replace("|", "/")
                        for cell in row.cells
                    ]
                    if not push("| " + " | ".join(cells) + " |"):
                        return "".join(out)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            if not push("\n발표자 노트:"):
                return "".join(out)
            if not push(notes):
                return "".join(out)
        push("")
    return "".join(out).strip()


def rtf_to_markdown(data: bytes, max_chars: int) -> str:
    from striprtf.striprtf import rtf_to_text

    text = data.decode("latin-1", errors="replace") if isinstance(data, bytes) else str(data)
    return rtf_to_text(text, errors="ignore").strip()[:max_chars]


def html_attachment_to_markdown(data: bytes, max_chars: int) -> str:
    text = _decode_text(data)
    try:
        md = html_to_md(text, heading_style="ATX")
    except Exception as exc:
        log.warning("html attachment markdownify failed: %s", exc)
        md = text
    return md.strip()[:max_chars]


def text_to_markdown(data: bytes, max_chars: int, lang_hint: str = "") -> str:
    text = _decode_text(data)
    text = text.replace("\r\n", "\n").replace("\r", "\n").rstrip()
    if lang_hint:
        # leave room for the fences in max_chars budget
        budget = max(0, max_chars - len(lang_hint) - 8)
        return f"```{lang_hint}\n{text[:budget]}\n```"
    return text[:max_chars]


def _decode_text(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr", "cp1252", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _image_stub(filename: str, size: int) -> str:
    return f"(이미지 파일 — 본문 추출 미지원. 원본은 attachments/{filename}, {size:,} bytes)"


def _unsupported_stub(filename: str, size: int) -> str:
    suffix = Path(filename).suffix or "?"
    return f"(형식 {suffix} 미지원 — 원본은 attachments/{filename}, {size:,} bytes)"


# -----------------------------------------------------------------------------
# Master dispatch + unified document
# -----------------------------------------------------------------------------


def attachment_to_markdown(filename: str, data: bytes, max_chars: int) -> str:
    """Return markdown text for an attachment.

    Always returns a non-empty string so the caller can include the attachment
    in the unified body.md (with a stub for unsupported / image types). This
    keeps body.md a complete index of every attachment the mail had.
    """
    suffix = Path(filename).suffix.lower()
    size = len(data) if data else 0

    if not data:
        return _unsupported_stub(filename, 0)

    try:
        if suffix == ".pdf":
            md = pdf_to_markdown(data, max_chars)
        elif suffix == ".docx":
            md = docx_to_markdown(data, max_chars)
        elif suffix == ".xlsx":
            md = xlsx_to_markdown(data, max_chars)
        elif suffix == ".pptx":
            md = pptx_to_markdown(data, max_chars)
        elif suffix == ".rtf":
            md = rtf_to_markdown(data, max_chars)
        elif suffix in (".html", ".htm"):
            md = html_attachment_to_markdown(data, max_chars)
        elif suffix in _TEXT_SUFFIXES:
            md = text_to_markdown(data, max_chars, lang_hint=_TEXT_SUFFIXES[suffix])
        elif suffix in _IMAGE_SUFFIXES:
            return _image_stub(filename, size)
        else:
            return _unsupported_stub(filename, size)
    except Exception as exc:
        log.warning("failed to parse attachment %s (%s): %s", filename, suffix, exc)
        return _unsupported_stub(filename, size)

    md = (md or "").strip()
    return md if md else _unsupported_stub(filename, size)


def build_unified_markdown(
    subject: str,
    sender: str,
    received: str,
    folder_path: str,
    body_md: str,
    attachment_sections: list[tuple[str, str]],
) -> str:
    """Combine mail body + attachment markdown into one document."""
    header = (
        f"# {subject or '(제목 없음)'}\n\n"
        f"- **From:** {sender}\n"
        f"- **Received:** {received}\n"
        f"- **Folder:** {folder_path}\n\n"
        "---\n\n"
    )
    parts = [header, body_md.strip()]
    for filename, md in attachment_sections:
        parts.append(ATTACHMENT_HEADER.format(name=filename))
        parts.append(md.strip())
    return "\n".join(parts).strip() + "\n"
