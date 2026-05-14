from __future__ import annotations
import io
import re
import base64
from pathlib import Path
from typing import Any, List
from haystack import component, Document

try:
    from markdownify import markdownify as html_to_md
except ImportError:
    html_to_md = lambda html, **kwargs: html

_HTML_TAG_RE = re.compile(r"<[^>]+>")
_HTML_ENTITY_RE = re.compile(r"&(?:amp|lt|gt|quot|#39|nbsp);")
_HTML_ENTITIES = {"&amp;": "&", "&lt;": "<", "&gt;": ">", "&quot;": '"', "&#39;": "'", "&nbsp;": " "}

_TEXT_SUFFIXES = {
    ".txt": "", ".log": "", ".md": "", ".csv": "csv", ".tsv": "tsv", ".json": "json",
    ".xml": "xml", ".yaml": "yaml", ".yml": "yaml", ".ini": "ini", ".conf": "",
    ".cfg": "", ".sql": "sql", ".py": "python", ".js": "javascript", ".ts": "typescript",
    ".sh": "bash", ".bat": "bat", ".ps1": "powershell"
}

@component
class EmailMarkdownConverter:
    """Converts raw email data into unified Markdown Documents."""

    def __init__(self, max_body_chars: int = 1_000_000, max_attachment_chars: int = 500_000):
        self.max_body_chars = max_body_chars
        self.max_attachment_chars = max_attachment_chars

    @component.output_types(documents=List[Document])
    def run(self, raw_emails: List[dict[str, Any]]) -> dict[str, Any]:
        docs = []
        for raw in raw_emails:
            body_md = self._body_to_markdown(raw.get("body_html", ""), raw.get("body_plain", ""))
            
            attachment_sections = []
            for att in raw.get("attachments", []):
                data = base64.b64decode(att["data_b64"]) if "data_b64" in att else b""
                md = self._attachment_to_markdown(att["filename"], data)
                attachment_sections.append((att["filename"], md))
            
            unified = self._build_unified_markdown(
                subject=raw.get("subject", ""),
                sender=raw.get("sender", ""),
                received=raw.get("received", ""),
                folder_path=raw.get("folder_path", ""),
                body_md=body_md,
                attachment_sections=attachment_sections
            )
            
            # Create Document with the unified markdown and email metadata
            email_meta = {
                "kind": raw.get("kind", "email"),
                "source": raw.get("source", "eml"),
                "subject": raw.get("subject", ""),
                "sender": raw.get("sender", ""),
                "received": raw.get("received", ""),
                "folder_path": raw.get("folder_path", ""),
                "mail_id": raw.get("mail_id", "")
            }
            # Include any other fields in raw that are not binary
            for k, v in raw.items():
                if k not in email_meta and k not in ("body_html", "body_plain", "attachments"):
                    email_meta[k] = v

            docs.append(Document(content=unified, meta=email_meta))
        return {"documents": docs}

    def _body_to_markdown(self, html: str, plain: str) -> str:
        if html and html.strip():
            try: return html_to_md(html, heading_style="ATX").strip()[:self.max_body_chars]
            except: return (plain or "")[:self.max_body_chars]
        return (plain or "")[:self.max_body_chars]

    def _attachment_to_markdown(self, filename: str, data: bytes) -> str:
        suffix = Path(filename).suffix.lower()
        if not data: return f"(Attachment {filename} is empty)"
        
        try:
            if suffix == ".pdf": return self._pdf_to_md(data)
            if suffix == ".docx": return self._docx_to_md(data)
            if suffix == ".xlsx": return self._xlsx_to_md(data)
            if suffix == ".pptx": return self._pptx_to_md(data)
            if suffix in _TEXT_SUFFIXES: return self._text_to_md(data, _TEXT_SUFFIXES[suffix])
            return f"(Attachment {filename} [format {suffix}] is supported only as raw file)"
        except Exception as e:
            return f"(Failed to parse attachment {filename}: {e})"

    def _pdf_to_md(self, data: bytes) -> str:
        import fitz
        parts = []
        used = 0
        with fitz.open(stream=data, filetype="pdf") as doc:
            for i, page in enumerate(doc, 1):
                text = (page.get_text("text") or "").strip()
                if not text: continue
                chunk = f"### Page {i}\n\n{text}\n"
                if used + len(chunk) > self.max_attachment_chars: break
                parts.append(chunk)
                used += len(chunk)
        return "\n".join(parts)

    def _docx_to_md(self, data: bytes) -> str:
        from docx import Document as DocxDocument
        doc = DocxDocument(io.BytesIO(data))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip(): parts.append(p.text.strip())
        return "\n\n".join(parts)[:self.max_attachment_chars]

    def _xlsx_to_md(self, data: bytes) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"### Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                line = "| " + " | ".join(str(c).replace("\n", " ") if c is not None else "" for c in row) + " |"
                if len("\n".join(lines) + line) > self.max_attachment_chars: break
                lines.append(line)
        return "\n".join(lines)

    def _pptx_to_md(self, data: bytes) -> str:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"### Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts)[:self.max_attachment_chars]

    def _text_to_md(self, data: bytes, lang: str) -> str:
        text = data.decode("utf-8", errors="replace")[:self.max_attachment_chars]
        return f"```{lang}\n{text}\n```" if lang else text

    def _build_unified_markdown(self, subject, sender, received, folder_path, body_md, attachment_sections) -> str:
        header = f"# {subject or '(No Subject)'}\n\n- **From:** {sender}\n- **Received:** {received}\n- **Folder:** {folder_path}\n\n---\n\n"
        parts = [header, body_md]
        for filename, md in attachment_sections:
            parts.append(f"\n\n---\n\n[Attachment: {filename}]\n\n" + md)
        return "\n".join(parts)
